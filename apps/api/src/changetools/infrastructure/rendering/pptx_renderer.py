"""BrandedDeckRenderer — builds a Brief-driven .pptx using python-pptx.

Layout philosophy
-----------------
We don't try to clone PowerPoint's master/layout system. Instead we use a
single blank layout and place rectangles, text frames, and tables ourselves.
This gives us pixel-level control over brand application (fills, fonts,
accent stripe, logo placement, per-slide source footnote) at the cost of
re-implementing layout primitives. For a Brief-shaped deck that's a fair
trade — every slide has a predictable header/body/footer.

Slide order (16:9, 13.33in by 7.5in):

    1. Title — project name, brief headline, accent stripe
    2. Executive summary — body text + sources
    3. Themes — 1 slide listing the themes (or one slide per theme if many)
    4. KPIs — table: KPI / Current / Target / Trend
    5. Risks — table: ID / Title / Severity / Owner
    6. Stakeholders — table: Name / Role / Sentiment / Concern
    7. Recommendations — numbered list with priority badges
    8. Conflicts (if any) — disagreement summary

Citations are surfaced as a "Sources" footer per slide listing the unique
documents the slide draws on.
"""

from __future__ import annotations

import contextlib
import io
import logging
from collections.abc import Iterable, Sequence
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from changetools.domain.brand import BrandProfile, Neutrals
from changetools.domain.brief import (
    KPI,
    Brief,
    Confidence,
    Conflict,
    ExecutiveSummary,
    Recommendation,
    Risk,
    Stakeholder,
    Theme,
)
from changetools.domain.chunk import ChunkRef

log = logging.getLogger(__name__)

# Default to Obsidian palette so the renderer is usable even when no brand
# profile exists yet.
DEFAULT_PRIMARY = "#0D171E"
DEFAULT_SECONDARY = "#2C343C"
DEFAULT_ACCENT = "#3A8C91"
DEFAULT_NEUTRALS = {"stone": "#6B7280", "hairline": "#D9DEE3", "fog": "#EEF1F3"}
DEFAULT_HEADING = "Aptos Display"
DEFAULT_BODY = "Aptos"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Header band reserved for accent stripe + title block.
HEADER_H = Inches(0.85)
# Footer band for sources + page number + logo.
FOOTER_H = Inches(0.55)
# Margins for the body content.
LEFT_MARGIN = Inches(0.5)
RIGHT_MARGIN = Inches(0.5)
TOP_MARGIN_BODY = HEADER_H + Inches(0.3)
BOTTOM_MARGIN_BODY = FOOTER_H + Inches(0.2)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Bad hex colour: {hex_color}")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class BrandedDeckRenderer:
    """Render a ``Brief`` to a python-pptx ``Presentation``.

    Inputs that vary per-tenant (palette, fonts, logo bytes, project name)
    are passed in via the constructor; ``render(brief)`` is otherwise pure.
    """

    def __init__(
        self,
        *,
        brand: BrandProfile | None,
        project_name: str,
        logo_bytes: bytes | None,
    ) -> None:
        self.project_name = project_name
        self.logo_bytes = logo_bytes
        if brand:
            self.primary = _hex_to_rgb(brand.primary_color)
            self.secondary = _hex_to_rgb(brand.secondary_color)
            self.accent = _hex_to_rgb(brand.accent_color)
            self.stone = _hex_to_rgb(brand.neutrals.stone)
            self.hairline = _hex_to_rgb(brand.neutrals.hairline)
            self.fog = _hex_to_rgb(brand.neutrals.fog)
            self.font_heading = brand.font_heading
            self.font_body = brand.font_body
        else:
            self.primary = _hex_to_rgb(DEFAULT_PRIMARY)
            self.secondary = _hex_to_rgb(DEFAULT_SECONDARY)
            self.accent = _hex_to_rgb(DEFAULT_ACCENT)
            self.stone = _hex_to_rgb(DEFAULT_NEUTRALS["stone"])
            self.hairline = _hex_to_rgb(DEFAULT_NEUTRALS["hairline"])
            self.fog = _hex_to_rgb(DEFAULT_NEUTRALS["fog"])
            self.font_heading = DEFAULT_HEADING
            self.font_body = DEFAULT_BODY

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, brief: Brief) -> tuple[bytes, int]:
        """Render to .pptx bytes. Returns ``(pptx_bytes, slide_count)``."""
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        c = brief.content

        self._title_slide(prs, brief)
        self._executive_summary_slide(prs, c.executive_summary)
        if c.themes:
            self._themes_slide(prs, c.themes)
        if c.kpis:
            self._kpis_slide(prs, c.kpis)
        if c.risks:
            self._risks_slide(prs, c.risks)
        if c.stakeholders:
            self._stakeholders_slide(prs, c.stakeholders)
        if c.recommendations:
            self._recommendations_slide(prs, c.recommendations)
        if c.conflicts:
            self._conflicts_slide(prs, c.conflicts)

        slide_count = len(prs.slides)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), slide_count

    # ------------------------------------------------------------------
    # Slide constructors
    # ------------------------------------------------------------------

    def _title_slide(self, prs: Presentation, brief: Brief) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Full-bleed primary background.
        self._add_fill_rect(slide, 0, 0, SLIDE_W, SLIDE_H, self.primary)

        # Accent stripe across the top quarter.
        self._add_fill_rect(
            slide, 0, Inches(2.4), Inches(0.6), Inches(0.12), self.accent
        )

        # Logo top-left, if available, on title slide.
        self._add_logo(slide, x=Inches(0.6), y=Inches(0.55), height=Inches(0.7))

        # Project name (small, kerned).
        self._add_text(
            slide,
            text=self.project_name.upper(),
            x=Inches(0.6),
            y=Inches(2.0),
            width=Inches(12),
            height=Inches(0.4),
            font=self.font_body,
            size=Pt(12),
            bold=True,
            color=self.accent,
            spacing=4,
        )
        # Headline.
        self._add_text(
            slide,
            text=brief.content.executive_summary.headline,
            x=Inches(0.6),
            y=Inches(2.7),
            width=Inches(12),
            height=Inches(2.5),
            font=self.font_heading,
            size=Pt(40),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        # Subtitle.
        self._add_text(
            slide,
            text=f"Change strategy brief — v{brief.version}",
            x=Inches(0.6),
            y=Inches(5.6),
            width=Inches(10),
            height=Inches(0.4),
            font=self.font_body,
            size=Pt(14),
            color=self.fog,
        )

    def _executive_summary_slide(
        self, prs: Presentation, summary: ExecutiveSummary
    ) -> None:
        slide = self._content_slide(prs, "Executive summary")
        self._add_text(
            slide,
            text=summary.body,
            x=LEFT_MARGIN,
            y=TOP_MARGIN_BODY,
            width=SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN,
            height=Inches(4.0),
            font=self.font_body,
            size=Pt(18),
            color=self.secondary,
            line_spacing=1.4,
        )
        self._add_confidence_badge(slide, summary.confidence)
        self._add_sources_footer(slide, summary.evidence)

    def _themes_slide(self, prs: Presentation, themes: Sequence[Theme]) -> None:
        slide = self._content_slide(prs, "Themes")
        col_w = (SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN) / 2
        gap = Inches(0.25)
        # Render up to 6 themes in a 2-column grid.
        for i, theme in enumerate(list(themes)[:6]):
            row = i // 2
            col = i % 2
            x = LEFT_MARGIN + col * col_w
            y = TOP_MARGIN_BODY + row * (Inches(1.6) + gap)
            box_w = col_w - gap
            box_h = Inches(1.6)
            # Card.
            self._add_fill_rect(slide, x, y, box_w, box_h, self.fog)
            # Accent bar on left.
            self._add_fill_rect(slide, x, y, Inches(0.06), box_h, self.accent)
            # Title.
            self._add_text(
                slide,
                text=theme.title,
                x=x + Inches(0.18),
                y=y + Inches(0.15),
                width=box_w - Inches(0.3),
                height=Inches(0.4),
                font=self.font_heading,
                size=Pt(16),
                bold=True,
                color=self.primary,
            )
            # Body.
            self._add_text(
                slide,
                text=theme.description,
                x=x + Inches(0.18),
                y=y + Inches(0.55),
                width=box_w - Inches(0.3),
                height=Inches(1.0),
                font=self.font_body,
                size=Pt(11),
                color=self.secondary,
            )
        self._add_sources_footer(slide, _evidence_union(themes))

    def _kpis_slide(self, prs: Presentation, kpis: Sequence[KPI]) -> None:
        slide = self._content_slide(prs, "Key performance indicators")
        rows = [["KPI", "Current", "Target", "Trend"]]
        for k in kpis[:8]:
            rows.append(
                [
                    k.name,
                    k.current_value,
                    k.target_value or "—",
                    k.trend or "—",
                ]
            )
        self._add_table(slide, rows, col_widths=(0.42, 0.18, 0.20, 0.20))
        self._add_sources_footer(slide, _evidence_union(kpis))

    def _risks_slide(self, prs: Presentation, risks: Sequence[Risk]) -> None:
        slide = self._content_slide(prs, "Risks")
        rows = [["ID", "Title", "Severity", "Owner"]]
        for r in risks[:8]:
            rows.append([r.risk_id, r.title, r.severity, r.owner or "—"])
        self._add_table(slide, rows, col_widths=(0.10, 0.55, 0.15, 0.20))
        self._add_sources_footer(slide, _evidence_union(risks))

    def _stakeholders_slide(
        self, prs: Presentation, stakeholders: Sequence[Stakeholder]
    ) -> None:
        slide = self._content_slide(prs, "Stakeholders")
        rows = [["Name", "Role", "Sentiment", "Concern"]]
        for s in stakeholders[:7]:
            rows.append([s.name, s.role, s.sentiment, _truncate(s.concern, 90)])
        self._add_table(slide, rows, col_widths=(0.18, 0.18, 0.14, 0.50))
        self._add_sources_footer(slide, _evidence_union(stakeholders))

    def _recommendations_slide(
        self, prs: Presentation, recs: Sequence[Recommendation]
    ) -> None:
        slide = self._content_slide(prs, "Recommendations")
        body_y = TOP_MARGIN_BODY
        for _, rec in enumerate(recs[:6]):
            row_h = Inches(0.85)
            # Priority badge.
            self._add_fill_rect(slide, LEFT_MARGIN, body_y, Inches(0.6), row_h, self.accent)
            self._add_text(
                slide,
                text=rec.priority,
                x=LEFT_MARGIN,
                y=body_y,
                width=Inches(0.6),
                height=row_h,
                font=self.font_heading,
                size=Pt(16),
                bold=True,
                color=RGBColor(255, 255, 255),
                align="center",
                anchor="middle",
            )
            # Title + body.
            x = LEFT_MARGIN + Inches(0.85)
            w = SLIDE_W - x - RIGHT_MARGIN
            self._add_text(
                slide,
                text=rec.title,
                x=x,
                y=body_y + Inches(0.05),
                width=w,
                height=Inches(0.35),
                font=self.font_heading,
                size=Pt(14),
                bold=True,
                color=self.primary,
            )
            self._add_text(
                slide,
                text=rec.description,
                x=x,
                y=body_y + Inches(0.4),
                width=w,
                height=Inches(0.4),
                font=self.font_body,
                size=Pt(10),
                color=self.secondary,
            )
            body_y += row_h + Inches(0.12)
            if body_y > SLIDE_H - FOOTER_H - Inches(0.6):
                break
        self._add_sources_footer(slide, _evidence_union(recs))

    def _conflicts_slide(self, prs: Presentation, conflicts: Sequence[Conflict]) -> None:
        slide = self._content_slide(prs, "Source conflicts")
        rows = [["Topic", "Call notes view", "Data pack view"]]
        for c in conflicts[:5]:
            rows.append([c.topic, _truncate(c.position_a, 90), _truncate(c.position_b, 90)])
        self._add_table(slide, rows, col_widths=(0.30, 0.35, 0.35))
        # Conflicts have bifurcated sources — flatten both sides.
        all_refs = []
        for c in conflicts:
            all_refs.extend(c.sources_a)
            all_refs.extend(c.sources_b)
        self._add_sources_footer(slide, all_refs)

    # ------------------------------------------------------------------
    # Slide chrome (header bar, footer, logo, etc.)
    # ------------------------------------------------------------------

    def _content_slide(self, prs: Presentation, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Header band.
        self._add_fill_rect(slide, 0, 0, SLIDE_W, HEADER_H, self.primary)
        # Accent underline.
        self._add_fill_rect(
            slide, LEFT_MARGIN, HEADER_H, Inches(1.0), Inches(0.06), self.accent
        )
        # Title.
        self._add_text(
            slide,
            text=title,
            x=LEFT_MARGIN,
            y=Inches(0.18),
            width=Inches(11),
            height=Inches(0.5),
            font=self.font_heading,
            size=Pt(22),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        # Logo top-right.
        self._add_logo(
            slide,
            x=SLIDE_W - Inches(1.4),
            y=Inches(0.15),
            height=Inches(0.55),
        )
        return slide

    def _add_sources_footer(
        self, slide, evidence: Iterable[ChunkRef]
    ) -> None:
        unique = []
        seen: set[str] = set()
        for ev in evidence:
            name = ev.document_filename or ""
            if name and name not in seen:
                seen.add(name)
                unique.append(name)
        text = "Sources: " + ", ".join(unique) if unique else ""
        # Footer line.
        self._add_fill_rect(
            slide,
            0,
            SLIDE_H - FOOTER_H,
            SLIDE_W,
            Inches(0.04),
            self.hairline,
        )
        if text:
            self._add_text(
                slide,
                text=text,
                x=LEFT_MARGIN,
                y=SLIDE_H - FOOTER_H + Inches(0.1),
                width=SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN,
                height=Inches(0.4),
                font=self.font_body,
                size=Pt(9),
                color=self.stone,
            )

    def _add_confidence_badge(self, slide, confidence: Confidence) -> None:
        labels = {
            "high": ("HIGH CONFIDENCE", self.accent),
            "medium": ("MEDIUM CONFIDENCE", self.stone),
            "low": ("LOW CONFIDENCE", _hex_to_rgb("#B45309")),
        }
        label, color = labels[confidence]
        x = SLIDE_W - RIGHT_MARGIN - Inches(2.2)
        y = TOP_MARGIN_BODY - Inches(0.05)
        self._add_fill_rect(slide, x, y, Inches(2.0), Inches(0.32), color)
        self._add_text(
            slide,
            text=label,
            x=x,
            y=y,
            width=Inches(2.0),
            height=Inches(0.32),
            font=self.font_body,
            size=Pt(9),
            bold=True,
            color=RGBColor(255, 255, 255),
            align="center",
            anchor="middle",
            spacing=2,
        )

    # ------------------------------------------------------------------
    # Low-level shape helpers
    # ------------------------------------------------------------------

    def _add_fill_rect(self, slide, x, y, w, h, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    def _add_text(
        self,
        slide,
        *,
        text: str,
        x: Emu | int,
        y: Emu | int,
        width: Emu | int,
        height: Emu | int,
        font: str,
        size: Pt,
        bold: bool = False,
        color: RGBColor | None = None,
        align: str = "left",
        anchor: str = "top",
        spacing: int | None = None,
        line_spacing: float | None = None,
    ) -> None:
        tx = slide.shapes.add_textbox(x, y, width, height)
        tx.text_frame.word_wrap = True
        tx.text_frame.margin_left = 0
        tx.text_frame.margin_right = 0
        tx.text_frame.margin_top = 0
        tx.text_frame.margin_bottom = 0
        if anchor == "middle":
            tx.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tx.text_frame.paragraphs[0]
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        if line_spacing is not None:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = size
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        if spacing is not None:
            # python-pptx doesn't expose letter-spacing directly; we approximate
            # by widening slightly via spc (1/100 pt). Best-effort.
            with contextlib.suppress(Exception):
                run.font._rPr.set("spc", str(spacing * 100))

    def _add_logo(self, slide, *, x, y, height) -> None:
        if not self.logo_bytes:
            return
        try:
            slide.shapes.add_picture(io.BytesIO(self.logo_bytes), x, y, height=height)
        except Exception as exc:  # pragma: no cover - defensive
            log.warning("deck.logo_embed_failed", exc_info=exc)

    def _add_table(
        self,
        slide,
        rows: Sequence[Sequence[Any]],
        *,
        col_widths: Sequence[float],
    ) -> None:
        """Add a table that fills the body area.

        ``col_widths`` are fractions of the body width (must sum to 1.0).
        """
        if not rows:
            return
        body_x = LEFT_MARGIN
        body_y = TOP_MARGIN_BODY
        body_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
        body_h = SLIDE_H - TOP_MARGIN_BODY - BOTTOM_MARGIN_BODY

        n_rows = len(rows)
        n_cols = len(rows[0])
        table_shape = slide.shapes.add_table(n_rows, n_cols, body_x, body_y, body_w, body_h)
        table = table_shape.table

        # Column widths.
        total = sum(col_widths)
        for i, frac in enumerate(col_widths):
            table.columns[i].width = Emu(int(body_w * (frac / total)))

        # Header row styling.
        for j, header in enumerate(rows[0]):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.primary
            self._format_cell(cell, str(header), color=RGBColor(255, 255, 255), bold=True)

        # Body rows.
        for i, row in enumerate(rows[1:], start=1):
            for j, value in enumerate(row):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.fog if i % 2 == 0 else RGBColor(255, 255, 255)
                self._format_cell(cell, str(value), color=self.secondary)

    def _format_cell(
        self,
        cell,
        text: str,
        *,
        color: RGBColor,
        bold: bool = False,
    ) -> None:
        cell.text = ""
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = self.font_body if not bold else self.font_heading
        run.font.size = Pt(11)
        run.font.bold = bold
        run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _evidence_union(items: Iterable[Any]) -> list[ChunkRef]:
    """Collect ChunkRef references across a list of items that each have an ``evidence`` field."""
    out: list[ChunkRef] = []
    for it in items:
        evs = getattr(it, "evidence", None) or []
        out.extend(evs)
    return out


def _truncate(s: str, n: int) -> str:
    return s if len(s) <= n else s[: n - 1].rstrip() + "…"


# ``Neutrals`` is referenced indirectly through ``BrandProfile`` — keep the
# import alive for type-checkers reading public re-exports.
_: type[Neutrals] = Neutrals
